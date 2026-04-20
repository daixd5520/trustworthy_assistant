import base64
import csv
import json
import os
import shlex
import subprocess
import urllib.request
from dataclasses import dataclass
from datetime import datetime
from pathlib import Path
from typing import TYPE_CHECKING, Any, Callable
from openai import OpenAI

from trustworthy_assistant.bookkeeping import BookkeepingService
from trustworthy_assistant.memory.service import TrustworthyMemoryService
from trustworthy_assistant.supervisor.models import ReviewFinding, Severity, TaskPhase

if TYPE_CHECKING:
    from trustworthy_assistant.supervisor.workflow import SupervisorWorkflow

_HOME = Path.home()
_BLOCKED_PREFIXES = [
    Path("/etc"), Path("/var"), Path("/sys"), Path("/proc"), Path("/dev"),
    Path("/sbin"), Path("/usr/sbin"),
    _HOME / ".ssh",
    _HOME / ".gnupg",
    _HOME / ".aws",
    _HOME / ".kube",
]
_COMMAND_TIMEOUT_DEFAULT = 20


def _truncate_text(text: str, max_chars: int) -> str:
    if len(text) <= max_chars:
        return text
    return text[:max_chars] + f"\n\n[... truncated, total {len(text)} chars ...]"


def _looks_binary(data: bytes) -> bool:
    if not data:
        return False
    if b"\x00" in data:
        return True
    sample = data[:1024]
    text_bytes = sum(1 for byte in sample if byte in b"\t\n\r" or 32 <= byte <= 126)
    return text_bytes / max(1, len(sample)) < 0.75
_COMMAND_TIMEOUT_MAX = 30
_COMMAND_OUTPUT_DEFAULT = 4000
_COMMAND_OUTPUT_MAX = 12000
_SHELL_META_TOKENS = ["&&", "||", ";", "|", ">", "<", "$(", "`", "\n"]
_ALWAYS_BLOCKED_COMMANDS = {
    "sudo", "su", "rm", "dd", "mkfs", "diskutil", "shutdown", "reboot",
    "launchctl", "chmod", "chown", "kill", "killall", "pkill", "scp", "ssh",
    "curl", "wget", "nc", "ncat", "telnet",
}
_BLOCKED_GIT_SUBCOMMANDS = {
    "reset", "checkout", "switch", "restore", "clean", "rebase", "merge",
    "cherry-pick", "push", "pull", "commit", "stash", "tag", "am", "apply",
}
_MEDIUM_RISK_GIT_SUBCOMMANDS = {"status", "diff", "log", "show", "branch"}
_MEDIUM_RISK_NODE_SUBCOMMANDS = {"test", "run", "exec"}
_BLOCKED_NODE_SUBCOMMANDS = {"install", "add", "remove", "update", "upgrade", "publish"}
_CONFIRMATION_HINTS = [
    "确认执行", "批准执行", "允许执行", "可以执行", "执行这个命令",
    "approve", "approved", "run it", "run this command", "execute it",
]


def _debug_emit(hypothesis_id: str, location: str, msg: str, data: dict[str, Any] | None = None, trace_id: str = "") -> None:
    payload = {
        "sessionId": "vision-read-image",
        "runId": "pre-fix",
        "hypothesisId": hypothesis_id,
        "location": location,
        "msg": f"[DEBUG] {msg}",
        "data": data or {},
    }
    if trace_id:
        payload["traceId"] = trace_id
    try:
        urllib.request.urlopen(
            urllib.request.Request(
                "http://127.0.0.1:7778/event",
                data=json.dumps(payload, ensure_ascii=False).encode("utf-8"),
                headers={"Content-Type": "application/json"},
            ),
            timeout=1.5,
        ).read()
    except Exception:
        pass


def _is_minimax_text_model(model_id: str) -> bool:
    normalized = (model_id or "").strip()
    return normalized in {
        "MiniMax-M2.7",
        "MiniMax-M2.7-highspeed",
    }


def _is_minimax_vision_model(model_id: str) -> bool:
    normalized = (model_id or "").strip()
    return normalized in {
        "MiniMax-VL-01",
    }


def _derive_minimax_openai_base_url(base_url: str) -> str:
    normalized = (base_url or "").strip().rstrip("/")
    if not normalized:
        return "https://api.minimaxi.com/v1"
    if normalized.endswith("/anthropic"):
        return normalized[: -len("/anthropic")] + "/v1"
    if normalized.endswith("/anthropic/messages"):
        return normalized[: -len("/anthropic/messages")] + "/v1"
    return normalized


@dataclass(slots=True)
class PendingCommandApproval:
    session_key: str
    command: str
    argv: list[str]
    cwd: str
    risk: str
    reason: str
    timeout_seconds: int
    max_output_chars: int
    created_at: str


def _local_now_str() -> str:
    now = datetime.now().astimezone()
    offset = now.strftime("%z")
    offset_fmt = f"UTC{offset[:3]}:{offset[3:]}" if len(offset) == 5 else offset
    return now.strftime(f"%Y-%m-%d %H:%M {offset_fmt}")


def _resolve_safe_path(path: str, workspace_dir: Path) -> Path:
    raw = path.strip()
    if raw.startswith("~"):
        candidate = Path(raw).expanduser().resolve()
    elif raw.startswith("/"):
        candidate = Path(raw).resolve()
    else:
        candidate = (workspace_dir / raw).resolve()
    for blocked in _BLOCKED_PREFIXES:
        try:
            candidate.relative_to(blocked)
            raise ValueError(f"Access denied: {path} is inside a protected directory")
        except ValueError:
            if candidate == blocked or blocked in candidate.parents:
                raise ValueError(f"Access denied: {path} is inside a protected directory")
    try:
        candidate.relative_to(_HOME)
    except ValueError:
        if candidate != _HOME and _HOME not in candidate.parents:
            raise ValueError(f"Access denied: {path} is outside home directory")
    return candidate


class ToolRegistry:
    def __init__(
        self,
        memory_service: TrustworthyMemoryService,
        bookkeeping_service: BookkeepingService,
        on_tool: Callable[[str, str], None] | None = None,
        reminder_callback: Callable[[str, int, str, str], None] | None = None,
        file_sender: Callable[[str, str, str], None] | None = None,
        message_sender: Callable[[str, str, str], None] | None = None,
        anthropic_client: Any | None = None,
        anthropic_api_key: str = "",
        anthropic_base_url: str | None = None,
        model_id: str = "",
        vision_api_key: str | None = None,
        vision_base_url: str | None = None,
        vision_model_id: str | None = None,
        supervisor_workflow: "SupervisorWorkflow | None" = None,
        state_dir: Path | None = None,
    ) -> None:
        self.memory_service = memory_service
        self.bookkeeping_service = bookkeeping_service
        self.on_tool = on_tool
        self.reminder_callback = reminder_callback
        self.file_sender = file_sender
        self.message_sender = message_sender
        self.anthropic_client = anthropic_client
        self.anthropic_api_key = anthropic_api_key
        self.anthropic_base_url = anthropic_base_url
        self.model_id = model_id
        self.vision_api_key = vision_api_key
        self.vision_base_url = vision_base_url
        self.vision_model_id = vision_model_id
        if _is_minimax_vision_model(self.vision_model_id or ""):
            if not self.vision_api_key and self.anthropic_api_key:
                self.vision_api_key = self.anthropic_api_key
            if not self.vision_base_url:
                self.vision_base_url = _derive_minimax_openai_base_url(self.anthropic_base_url or "")
        self._vision_client: OpenAI | None = None
        self._minimax_vision_client: OpenAI | None = None
        self.supervisor_workflow = supervisor_workflow
        self.workspace_dir = self.memory_service.repository.paths.workspace_dir
        self._current_channel: str = "terminal"
        self._current_user_id: str = "local"
        self._current_agent_id: str = "main"
        self._latest_user_input: str = ""
        self._current_session_key: str = ""
        self._pending_command_approvals: dict[str, PendingCommandApproval] = {}
        self._approved_command_prefixes: dict[str, list[tuple[str, ...]]] = {}
        self._state_dir = state_dir or (self.workspace_dir / ".assistant_state")
        self._state_error: str = ""
        self._state_file: Path | None = None
        try:
            self._state_dir.mkdir(parents=True, exist_ok=True)
            self._state_file = self._state_dir / "tool_registry_state.json"
        except OSError as exc:
            self._state_error = f"Approval persistence disabled: {exc}"
            self.emit("run_command", self._state_error)
        self._load_state()
        self.tools = [
            {
                "name": "memory_write",
                "description": "Save an important fact or observation to long-term memory.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "content": {"type": "string", "description": "The fact or observation to remember."},
                        "category": {"type": "string", "description": "Category: preference, fact, context, etc."},
                    },
                    "required": ["content"],
                },
            },
            {
                "name": "memory_search",
                "description": "Search stored memories for relevant information, ranked by similarity.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "query": {"type": "string", "description": "Search query."},
                        "top_k": {"type": "integer", "description": "Max results. Default: 5."},
                    },
                    "required": ["query"],
                },
            },
            {
                "name": "ledger_add_entry",
                "description": "Add one bookkeeping entry to the local ledger for the current user. Use for expenses, income, reimbursements, and transfers that should be tracked.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "amount": {"type": "number", "description": "Positive amount, such as 32.5."},
                        "category": {"type": "string", "description": "Category like food, transport, rent, shopping, salary, or reimbursement."},
                        "entry_type": {"type": "string", "description": "Either `expense` or `income`. Default: `expense`."},
                        "note": {"type": "string", "description": "Optional note, merchant, or context for the entry."},
                        "occurred_at": {"type": "string", "description": "Optional ISO datetime for when the transaction happened. Default: now."},
                        "currency": {"type": "string", "description": "Currency code like CNY, USD. Default: CNY."},
                        "account": {"type": "string", "description": "Optional account name like cash, wechat, alipay, bank_card."},
                        "source": {"type": "string", "description": "Optional source label, such as wechat, manual, reimbursement."},
                    },
                    "required": ["amount", "category"],
                },
            },
            {
                "name": "ledger_report",
                "description": "Generate a bookkeeping report with totals and category stats. Use when the user asks for today's, this week's, last week's, this month's, or last month's账本/账单/统计.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "period": {"type": "string", "description": "One of: today, yesterday, week, last_week, month, last_month."},
                        "tz": {"type": "string", "description": "Optional timezone name. Default: Local."},
                    },
                    "required": ["period"],
                },
            },
            {
                "name": "ledger_configure_reports",
                "description": "Configure automatic ledger reports for the current user and channel. Use when the user asks for daily, weekly, or monthly账单推送.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "daily_enabled": {"type": "boolean", "description": "Whether to enable the daily report. Default: true."},
                        "weekly_enabled": {"type": "boolean", "description": "Whether to enable the weekly report. Default: true."},
                        "monthly_enabled": {"type": "boolean", "description": "Whether to enable the monthly report. Default: true."},
                        "daily_time": {"type": "string", "description": "HH:MM local time for the daily report. Default: 23:00."},
                        "weekly_time": {"type": "string", "description": "HH:MM local time for the weekly report. Default: 23:00."},
                        "monthly_time": {"type": "string", "description": "HH:MM local time for the monthly report. Default: 00:05."},
                        "weekly_weekday": {"type": "string", "description": "sun, mon, tue, wed, thu, fri, or sat. Default: sun."},
                        "tz": {"type": "string", "description": "Timezone name. Default: Local."},
                    },
                },
            },
            {
                "name": "list_directory",
                "description": "List files and directories. Supports workspace-relative paths, absolute paths, and ~/ paths. Access is limited to your home directory with system directories blocked.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to list. Can be relative (workspace), absolute (/Users/...), or home-relative (~/Downloads). Default: ."},
                    },
                },
            },
            {
                "name": "read_file",
                "description": "Read a text file so you can analyze it before answering. Supports workspace-relative paths, absolute paths, and ~/ paths. Access is limited to your home directory with system directories blocked. After reading, default to summarizing and evaluating the file instead of quoting its contents unless the user explicitly asks to see the content.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "File path. Can be relative (workspace), absolute (/Users/...), or home-relative (~/Documents/notes.txt)."},
                        "max_chars": {"type": "integer", "description": "Max characters to return. Default: 4000."},
                    },
                    "required": ["path"],
                },
            },
            {
                "name": "write_file",
                "description": "Create or overwrite a text file on the local filesystem. Use this instead of shell redirection when you need to write code, config, notes, or other text files.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Target file path. Supports workspace-relative, absolute, and ~/ paths."},
                        "content": {"type": "string", "description": "The full text content to write."},
                        "overwrite": {"type": "boolean", "description": "Whether to overwrite an existing file. Default: true."},
                        "create_parent_dirs": {"type": "boolean", "description": "Whether to create missing parent directories automatically. Default: true."},
                    },
                    "required": ["path", "content"],
                },
            },
            {
                "name": "append_file",
                "description": "Append text to a local file. Use for logs, notes, markdown, or incremental text output without relying on shell redirects.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Target file path. Supports workspace-relative, absolute, and ~/ paths."},
                        "content": {"type": "string", "description": "Text to append."},
                        "create_parent_dirs": {"type": "boolean", "description": "Whether to create missing parent directories automatically. Default: true."},
                    },
                    "required": ["path", "content"],
                },
            },
            {
                "name": "replace_in_file",
                "description": "Edit an existing text file by replacing a target string. Use this for focused updates without rewriting the whole file.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Existing file path. Supports workspace-relative, absolute, and ~/ paths."},
                        "old_text": {"type": "string", "description": "Exact text to find."},
                        "new_text": {"type": "string", "description": "Replacement text."},
                        "replace_all": {"type": "boolean", "description": "Whether to replace all matches. Default: false."},
                    },
                    "required": ["path", "old_text", "new_text"],
                },
            },
            {
                "name": "make_directory",
                "description": "Create a directory on the local filesystem, including parent directories if needed.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Directory path. Supports workspace-relative, absolute, and ~/ paths."},
                    },
                    "required": ["path"],
                },
            },
            {
                "name": "get_current_time",
                "description": "Get the current local time with timezone.",
                "input_schema": {"type": "object", "properties": {}},
            },
            {
                "name": "set_reminder",
                "description": "Set a one-off reminder that fires after a delay. The reminder message will be sent to you when it triggers.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "message": {"type": "string", "description": "The reminder message to send when it triggers."},
                        "delay_minutes": {"type": "integer", "description": "Number of minutes from now to trigger the reminder."},
                    },
                    "required": ["message", "delay_minutes"],
                },
            },
            {
                "name": "send_file",
                "description": "Send a file from the local filesystem to the user via the current channel (e.g., WeChat). Supports images, videos, documents, and other files up to 20MB. Images and videos are sent in their native format for optimal display. Only works when connected through a channel that supports file delivery (like WeChat).",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to the file to send. Can be relative (workspace), absolute (/Users/...), or home-relative (~/Documents/report.pdf)."},
                        "caption": {"type": "string", "description": "Optional text caption to include with the file."},
                    },
                    "required": ["path"],
                },
            },
            {
                "name": "read_image",
                "description": "Read and describe a local image file using the vision model. Useful when the user sends a screenshot, photo, chart, or image attachment.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "path": {"type": "string", "description": "Path to the image file. Can be relative (workspace), absolute (/Users/...), or home-relative (~/Pictures/example.png)."},
                        "prompt": {"type": "string", "description": "Optional question or instruction for the image, such as 'Describe the screenshot and extract the key text'."},
                    },
                    "required": ["path"],
                },
            },
            {
                "name": "run_command",
                "description": "Run a terminal command under supervisor governance. Commands run without a shell, are limited to safe working directories, and dangerous commands are blocked. Low-risk read/test commands can run directly; higher-risk commands require explicit user confirmation in the latest user message.",
                "input_schema": {
                    "type": "object",
                    "properties": {
                        "command": {"type": "string", "description": "Terminal command to run, like `git status` or `pytest -q`. Shell operators like `&&`, `|`, `>` are not allowed."},
                        "cwd": {"type": "string", "description": "Optional working directory. Can be relative to the workspace, absolute under home, or ~/ path. Default: workspace root."},
                        "timeout_seconds": {"type": "integer", "description": "Execution timeout in seconds. Default: 20, max: 30."},
                        "max_output_chars": {"type": "integer", "description": "Maximum combined stdout/stderr to return. Default: 4000, max: 12000."},
                    },
                    "required": ["command"],
                },
            },
        ]
        self.handlers: dict[str, Callable[..., str]] = {
            "memory_write": self.memory_write,
            "memory_search": self.memory_search,
            "ledger_add_entry": self.ledger_add_entry,
            "ledger_report": self.ledger_report,
            "ledger_configure_reports": self.ledger_configure_reports,
            "list_directory": self.list_directory,
            "read_file": self.read_file,
            "write_file": self.write_file,
            "append_file": self.append_file,
            "replace_in_file": self.replace_in_file,
            "make_directory": self.make_directory,
            "get_current_time": self.get_current_time,
            "set_reminder": self.set_reminder,
            "send_file": self.send_file,
            "read_image": self.read_image,
            "run_command": self.run_command,
        }

    def _load_state(self) -> None:
        if self._state_file is None:
            return
        if not self._state_file.is_file():
            return
        try:
            payload = json.loads(self._state_file.read_text(encoding="utf-8"))
        except Exception as exc:
            self._state_error = f"Failed to load approval state: {exc}"
            self.emit("run_command", self._state_error)
            return
        pending_rows = payload.get("pending_command_approvals") if isinstance(payload, dict) else None
        prefix_rows = payload.get("approved_command_prefixes") if isinstance(payload, dict) else None
        if isinstance(pending_rows, dict):
            for session_key, row in pending_rows.items():
                if not isinstance(row, dict):
                    continue
                try:
                    self._pending_command_approvals[str(session_key)] = PendingCommandApproval(
                        session_key=str(row.get("session_key") or session_key),
                        command=str(row.get("command") or ""),
                        argv=[str(item) for item in row.get("argv", []) if str(item)],
                        cwd=str(row.get("cwd") or ""),
                        risk=str(row.get("risk") or "medium"),
                        reason=str(row.get("reason") or ""),
                        timeout_seconds=int(row.get("timeout_seconds") or _COMMAND_TIMEOUT_DEFAULT),
                        max_output_chars=int(row.get("max_output_chars") or _COMMAND_OUTPUT_DEFAULT),
                        created_at=str(row.get("created_at") or datetime.now().isoformat()),
                    )
                except Exception:
                    continue
        if isinstance(prefix_rows, dict):
            for session_key, rows in prefix_rows.items():
                if not isinstance(rows, list):
                    continue
                parsed_prefixes: list[tuple[str, ...]] = []
                for row in rows:
                    if not isinstance(row, list):
                        continue
                    prefix = tuple(str(item).lower() for item in row if str(item))
                    if prefix:
                        parsed_prefixes.append(prefix)
                if parsed_prefixes:
                    self._approved_command_prefixes[str(session_key)] = parsed_prefixes

    def _save_state(self) -> bool:
        if self._state_file is None:
            return False
        payload = {
            "pending_command_approvals": {
                session_key: {
                    "session_key": item.session_key,
                    "command": item.command,
                    "argv": item.argv,
                    "cwd": item.cwd,
                    "risk": item.risk,
                    "reason": item.reason,
                    "timeout_seconds": item.timeout_seconds,
                    "max_output_chars": item.max_output_chars,
                    "created_at": item.created_at,
                }
                for session_key, item in self._pending_command_approvals.items()
            },
            "approved_command_prefixes": {
                session_key: [list(prefix) for prefix in prefixes]
                for session_key, prefixes in self._approved_command_prefixes.items()
            },
        }
        try:
            self._state_file.parent.mkdir(parents=True, exist_ok=True)
            self._state_file.write_text(
                json.dumps(payload, ensure_ascii=False, indent=2) + "\n",
                encoding="utf-8",
            )
            self._state_error = ""
            return True
        except OSError as exc:
            self._state_error = f"Failed to persist approval state: {exc}"
            self.emit("run_command", self._state_error)
            return False

    def set_channel_context(
        self,
        channel: str,
        user_id: str,
        latest_user_input: str = "",
        session_key: str = "",
        agent_id: str = "main",
    ) -> None:
        self._current_channel = channel
        self._current_user_id = user_id
        self._current_agent_id = agent_id or "main"
        self._latest_user_input = latest_user_input
        self._current_session_key = session_key

    def emit(self, name: str, detail: str) -> None:
        if self.on_tool:
            self.on_tool(name, detail)

    def memory_write(self, content: str, category: str = "general") -> str:
        self.emit("memory_write", f"[{category}] {content[:60]}...")
        return self.memory_service.write_memory(
            content,
            category,
            agent_id=self._current_agent_id,
            channel=self._current_channel,
            user_id=self._current_user_id,
        )

    def memory_search(self, query: str, top_k: int = 5) -> str:
        self.emit("memory_search", query)
        results = self.memory_service.hybrid_search(
            query,
            top_k,
            agent_id=self._current_agent_id,
            channel=self._current_channel,
            user_id=self._current_user_id,
        )
        if not results:
            return "No relevant memories found."
        return "\n".join(
            f"[{item['path']}] (score: {item['score']}, status: {item['status']}) {item['snippet']}"
            for item in results
        )

    def ledger_add_entry(
        self,
        amount: float,
        category: str,
        entry_type: str = "expense",
        note: str = "",
        occurred_at: str = "",
        currency: str = "CNY",
        account: str = "cash",
        source: str = "manual",
    ) -> str:
        detail = f"{entry_type} {amount} {currency} {category}"
        self.emit("ledger_add_entry", detail)
        try:
            entry = self.bookkeeping_service.add_entry(
                amount=amount,
                category=category,
                entry_type=entry_type,
                note=note,
                occurred_at=occurred_at,
                currency=currency,
                account=account,
                source=source,
                channel=self._current_channel,
                user_id=self._current_user_id,
            )
        except Exception as exc:
            return f"Error: Failed to add ledger entry: {exc}"
        note_part = f" | note={entry.note}" if entry.note else ""
        return (
            f"Ledger entry saved: id={entry.entry_id} type={entry.entry_type} "
            f"amount={entry.amount} {entry.currency} category={entry.category} "
            f"occurred_at={entry.occurred_at}{note_part}"
        )

    def ledger_report(self, period: str, tz: str = "Local") -> str:
        self.emit("ledger_report", f"{period} tz={tz}")
        try:
            report = self.bookkeeping_service.summarize(period, tz_name=tz)
        except Exception as exc:
            return f"Error: Failed to build ledger report: {exc}"
        lines = [
            f"{report['label']}",
            f"Range: {report['start_at']} -> {report['end_at']}",
            f"Entries: {report['entry_count']}",
            f"Expense total: {report['expense_total']}",
            f"Income total: {report['income_total']}",
            f"Net total: {report['net_total']}",
        ]
        expense_rows = report.get("expense_by_category") or []
        income_rows = report.get("income_by_category") or []
        if expense_rows:
            lines.append("Expense by category:")
            lines.extend(f"- {row['category']}: {row['amount']}" for row in expense_rows[:8])
        if income_rows:
            lines.append("Income by category:")
            lines.extend(f"- {row['category']}: {row['amount']}" for row in income_rows[:8])
        entries = report.get("entries") or []
        if entries:
            lines.append("Recent entries:")
            for row in entries[:12]:
                note_part = f" | {row['note']}" if row.get("note") else ""
                lines.append(
                    f"- {row['occurred_at']} | {row['entry_type']} | {row['amount']} {row['currency']} | {row['category']}{note_part}"
                )
        else:
            lines.append("No ledger entries in this period.")
        return "\n".join(lines)

    def ledger_configure_reports(
        self,
        daily_enabled: bool = True,
        weekly_enabled: bool = True,
        monthly_enabled: bool = True,
        daily_time: str = "23:00",
        weekly_time: str = "23:00",
        monthly_time: str = "00:05",
        weekly_weekday: str = "sun",
        tz: str = "Local",
    ) -> str:
        self.emit(
            "ledger_configure_reports",
            (
                f"daily={daily_enabled}@{daily_time} weekly={weekly_enabled}@{weekly_time}/{weekly_weekday} "
                f"monthly={monthly_enabled}@{monthly_time} tz={tz}"
            ),
        )
        if not self._current_channel or not self._current_user_id:
            return "Error: Missing current channel or user context for scheduled delivery."
        try:
            jobs = self.bookkeeping_service.configure_report_jobs(
                channel=self._current_channel,
                user_id=self._current_user_id,
                tz_name=tz,
                daily_enabled=daily_enabled,
                weekly_enabled=weekly_enabled,
                monthly_enabled=monthly_enabled,
                daily_time=daily_time,
                weekly_time=weekly_time,
                monthly_time=monthly_time,
                weekly_weekday=weekly_weekday,
            )
        except Exception as exc:
            return f"Error: Failed to configure ledger reports: {exc}"
        lines = ["Ledger auto reports configured:"]
        for job in jobs:
            schedule = job.get("schedule") or {}
            lines.append(
                f"- {job.get('id')}: enabled={job.get('enabled')} expr={schedule.get('expr')} tz={schedule.get('tz')}"
            )
        lines.append("The scheduler will pick up the updated CRON.json automatically on the next reload cycle.")
        return "\n".join(lines)

    def list_directory(self, path: str = ".") -> str:
        self.emit("list_directory", path)
        target = _resolve_safe_path(path, self.workspace_dir)
        if not target.exists():
            return f"Error: Path not found: {path}"
        if not target.is_dir():
            return f"Error: Not a directory: {path}"
        rows = []
        for child in sorted(target.iterdir(), key=lambda item: (not item.is_dir(), item.name.lower()))[:200]:
            suffix = "/" if child.is_dir() else ""
            rows.append(child.name + suffix)
        return "\n".join(rows) if rows else "(empty)"

    def read_file(self, path: str, max_chars: int = 4000) -> str:
        self.emit("read_file", path)
        target = _resolve_safe_path(path, self.workspace_dir)
        if not target.exists():
            return f"Error: Path not found: {path}"
        if not target.is_file():
            return f"Error: Not a file: {path}"
        suffix = target.suffix.lower()
        if suffix == ".pdf":
            try:
                from pypdf import PdfReader
            except ImportError:
                return "Error: PDF reading backend is not installed."
            try:
                reader = PdfReader(str(target))
            except Exception as exc:
                return f"Error: Failed to open PDF: {exc}"
            parts: list[str] = []
            total_chars = 0
            for index, page in enumerate(reader.pages, start=1):
                try:
                    page_text = (page.extract_text() or "").strip()
                except Exception as exc:
                    page_text = f"[Page {index} extraction failed: {exc}]"
                if not page_text:
                    continue
                block = f"[Page {index}]\n{page_text}"
                parts.append(block)
                total_chars += len(block) + 2
                if total_chars >= max_chars:
                    break
            if not parts:
                return "Error: No extractable text found in PDF. It may be scanned, image-based, or encrypted."
            return _truncate_text("\n\n".join(parts), max_chars)
        if suffix == ".docx":
            try:
                from docx import Document
            except ImportError:
                return "Error: DOCX reading backend is not installed."
            try:
                document = Document(str(target))
            except Exception as exc:
                return f"Error: Failed to open DOCX: {exc}"
            parts: list[str] = []
            for paragraph in document.paragraphs:
                text = paragraph.text.strip()
                if text:
                    parts.append(text)
            for table_index, table in enumerate(document.tables, start=1):
                rows: list[str] = []
                for row in table.rows:
                    cells = [cell.text.strip() for cell in row.cells]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append(f"[Table {table_index}]\n" + "\n".join(rows))
            if not parts:
                return "Error: No extractable text found in DOCX."
            return _truncate_text("\n\n".join(parts), max_chars)
        if suffix == ".xlsx":
            try:
                from openpyxl import load_workbook
            except ImportError:
                return "Error: XLSX reading backend is not installed."
            try:
                workbook = load_workbook(filename=str(target), read_only=True, data_only=True)
            except Exception as exc:
                return f"Error: Failed to open XLSX: {exc}"
            parts: list[str] = []
            for sheet in workbook.worksheets[:5]:
                rows: list[str] = []
                for row in sheet.iter_rows(max_row=50, values_only=True):
                    cells = ["" if cell is None else str(cell).strip() for cell in row]
                    if any(cells):
                        rows.append(" | ".join(cells))
                if rows:
                    parts.append(f"[Sheet: {sheet.title}]\n" + "\n".join(rows))
            if not parts:
                return "Error: No extractable cells found in XLSX."
            return _truncate_text("\n\n".join(parts), max_chars)
        if suffix == ".pptx":
            try:
                from pptx import Presentation
            except ImportError:
                return "Error: PPTX reading backend is not installed."
            try:
                presentation = Presentation(str(target))
            except Exception as exc:
                return f"Error: Failed to open PPTX: {exc}"
            parts: list[str] = []
            for slide_index, slide in enumerate(list(presentation.slides)[:30], start=1):
                texts: list[str] = []
                for shape in slide.shapes:
                    text = getattr(shape, "text", "").strip()
                    if text:
                        texts.append(text)
                if texts:
                    parts.append(f"[Slide {slide_index}]\n" + "\n".join(texts))
            if not parts:
                return "Error: No extractable text found in PPTX."
            return _truncate_text("\n\n".join(parts), max_chars)
        if suffix == ".csv":
            try:
                with target.open("r", encoding="utf-8", errors="replace", newline="") as handle:
                    reader = csv.reader(handle)
                    rows = []
                    for index, row in enumerate(reader, start=1):
                        rows.append(" | ".join(cell.strip() for cell in row))
                        if index >= 100:
                            break
            except Exception as exc:
                return f"Error: Failed to open CSV: {exc}"
            if not rows:
                return "(empty csv)"
            return _truncate_text("\n".join(rows), max_chars)
        if suffix == ".json":
            try:
                parsed = json.loads(target.read_text(encoding="utf-8"))
            except Exception:
                parsed = None
            if parsed is not None:
                return _truncate_text(json.dumps(parsed, ensure_ascii=False, indent=2), max_chars)
        raw = target.read_bytes()
        if _looks_binary(raw):
            return (
                f"Error: Unsupported binary file type: {target.suffix or '(no extension)'}.\n"
                "Supported readable formats include: txt, md, json, csv, html, xml, pdf, docx, xlsx, pptx."
            )
        text = raw.decode("utf-8", errors="replace")
        return _truncate_text(text, max_chars)

    def write_file(
        self,
        path: str,
        content: str,
        overwrite: bool = True,
        create_parent_dirs: bool = True,
    ) -> str:
        self.emit("write_file", path)
        try:
            target = _resolve_safe_path(path, self.workspace_dir)
            parent = target.parent
            if create_parent_dirs:
                parent.mkdir(parents=True, exist_ok=True)
            elif not parent.is_dir():
                return f"Error: Parent directory does not exist: {parent}"
            if target.exists() and target.is_dir():
                return f"Error: Path is a directory: {target}"
            if target.exists() and not overwrite:
                return f"Error: File already exists: {target}"
            target.write_text(str(content), encoding="utf-8")
            return f"Wrote file: {target} ({len(str(content))} chars)"
        except Exception as exc:
            return f"Error: Failed to write file: {exc}"

    def append_file(
        self,
        path: str,
        content: str,
        create_parent_dirs: bool = True,
    ) -> str:
        self.emit("append_file", path)
        try:
            target = _resolve_safe_path(path, self.workspace_dir)
            parent = target.parent
            if create_parent_dirs:
                parent.mkdir(parents=True, exist_ok=True)
            elif not parent.is_dir():
                return f"Error: Parent directory does not exist: {parent}"
            if target.exists() and target.is_dir():
                return f"Error: Path is a directory: {target}"
            with target.open("a", encoding="utf-8") as handle:
                handle.write(str(content))
            return f"Appended to file: {target} (+{len(str(content))} chars)"
        except Exception as exc:
            return f"Error: Failed to append file: {exc}"

    def replace_in_file(
        self,
        path: str,
        old_text: str,
        new_text: str,
        replace_all: bool = False,
    ) -> str:
        self.emit("replace_in_file", path)
        if not old_text:
            return "Error: old_text must not be empty."
        try:
            target = _resolve_safe_path(path, self.workspace_dir)
            if not target.exists():
                return f"Error: File not found: {target}"
            if target.is_dir():
                return f"Error: Path is a directory: {target}"
            original = target.read_text(encoding="utf-8")
            if old_text not in original:
                return "Error: Target text not found in file."
            count = original.count(old_text) if replace_all else 1
            updated = original.replace(old_text, new_text) if replace_all else original.replace(old_text, new_text, 1)
            target.write_text(updated, encoding="utf-8")
            return f"Updated file: {target} ({count} replacement{'s' if count != 1 else ''})"
        except Exception as exc:
            return f"Error: Failed to replace text in file: {exc}"

    def make_directory(self, path: str) -> str:
        self.emit("make_directory", path)
        try:
            target = _resolve_safe_path(path, self.workspace_dir)
            target.mkdir(parents=True, exist_ok=True)
            return f"Directory ready: {target}"
        except Exception as exc:
            return f"Error: Failed to create directory: {exc}"

    def get_current_time(self) -> str:
        self.emit("get_current_time", "local")
        return _local_now_str()

    def set_reminder(self, message: str, delay_minutes: int) -> str:
        self.emit("set_reminder", f"{delay_minutes}m: {message[:40]}")
        if self.reminder_callback is None:
            return "Error: Reminders are not available in this session."
        if delay_minutes < 1:
            return "Error: delay_minutes must be at least 1."
        if delay_minutes > 1440:
            return "Error: delay_minutes must be at most 1440 (24 hours)."
        try:
            self.reminder_callback(message, delay_minutes, self._current_channel, self._current_user_id)
            return f"Reminder set: will notify you in {delay_minutes} minute(s)."
        except Exception as exc:
            return f"Error: Failed to set reminder: {exc}"

    def send_file(self, path: str, caption: str = "") -> str:
        detail = path if not caption else f"{path} | caption={caption[:40]}"
        self.emit("send_file", detail)
        if self.file_sender is None:
            return "Error: File sending is not available in this session. Only available when connected via WeChat."
        try:
            target = _resolve_safe_path(path, self.workspace_dir)
        except ValueError as exc:
            return f"Error: {exc}"
        if not target.exists():
            return f"Error: File not found: {path}"
        if not target.is_file():
            return f"Error: Not a file: {path}"
        file_size = target.stat().st_size
        if file_size > 20 * 1024 * 1024:
            return f"Error: File too large ({file_size} bytes). Maximum size is 20MB."
        if file_size == 0:
            return f"Error: File is empty: {path}"
        try:
            self.file_sender(str(target), self._current_channel, self._current_user_id)
            return f"File sent: {target.name} ({file_size} bytes)"
        except Exception as exc:
            return f"Error: Failed to send file: {exc}"

    @staticmethod
    def _image_media_type(path: Path, header: bytes = b"") -> str:
        suffix = path.suffix.lower()
        by_suffix = {
            ".png": "image/png",
            ".jpg": "image/jpeg",
            ".jpeg": "image/jpeg",
            ".gif": "image/gif",
            ".webp": "image/webp",
            ".bmp": "image/bmp",
        }
        if suffix in by_suffix:
            return by_suffix[suffix]
        if header.startswith(b"\x89PNG\r\n\x1a\n"):
            return "image/png"
        if header.startswith(b"\xff\xd8\xff"):
            return "image/jpeg"
        if header.startswith((b"GIF87a", b"GIF89a")):
            return "image/gif"
        if header.startswith(b"RIFF") and header[8:12] == b"WEBP":
            return "image/webp"
        if header.startswith(b"BM"):
            return "image/bmp"
        return ""

    def _should_use_mmx_vision(self) -> bool:
        minimax_base = (self.anthropic_base_url or "") + " " + (self.vision_base_url or "")
        return bool(
            self.anthropic_api_key
            and (
                _is_minimax_text_model(self.model_id)
                or _is_minimax_vision_model(self.vision_model_id or "")
                or "minimaxi.com" in minimax_base
            )
        )

    def _run_mmx_vision_describe(self, image_path: Path, prompt: str, trace_id: str = "") -> str:
        project_root = Path(__file__).resolve().parents[2]
        wrapper = project_root / ".dbg" / "mmx_from_env.py"
        if not wrapper.is_file():
            return ""
        command = [
            "python3",
            str(wrapper),
            "vision",
            "describe",
            "--image",
            str(image_path),
            "--prompt",
            prompt,
            "--output",
            "json",
            "--quiet",
        ]
        # #region debug-point V2:mmx-vision
        _debug_emit("V2", "tools.py:read_image", "trying mmx vision describe", {
            "command": "python3 .dbg/mmx_from_env.py vision describe",
            "image_path": str(image_path),
        }, trace_id=trace_id)
        # #endregion
        try:
            completed = subprocess.run(
                command,
                cwd=str(project_root),
                capture_output=True,
                text=True,
                timeout=120,
                check=False,
            )
        except Exception as exc:
            _debug_emit("V2", "tools.py:read_image", "mmx vision describe launch failed", {
                "error": str(exc)[:500],
            }, trace_id=trace_id)
            return ""
        stdout = (completed.stdout or "").strip()
        stderr = (completed.stderr or "").strip()
        if completed.returncode != 0:
            _debug_emit("V2", "tools.py:read_image", "mmx vision describe failed", {
                "returncode": completed.returncode,
                "stdout": stdout[:500],
                "stderr": stderr[:500],
            }, trace_id=trace_id)
            return ""
        json_start = stdout.find("{")
        json_end = stdout.rfind("}")
        if json_start == -1 or json_end == -1 or json_end < json_start:
            _debug_emit("V2", "tools.py:read_image", "mmx vision describe returned non-json output", {
                "stdout": stdout[:500],
            }, trace_id=trace_id)
            return ""
        try:
            payload = json.loads(stdout[json_start:json_end + 1])
        except Exception as exc:
            _debug_emit("V2", "tools.py:read_image", "mmx vision describe json parse failed", {
                "error": str(exc)[:300],
                "stdout": stdout[:500],
            }, trace_id=trace_id)
            return ""
        content = str(payload.get("content", "")).strip()
        if content:
            _debug_emit("V7", "tools.py:read_image", "mmx vision response summary", {
                "message_content_type": type(payload.get("content", "")).__name__,
                "base_status": payload.get("base_resp", {}).get("status_msg", ""),
            }, trace_id=trace_id)
        return content

    def read_image(self, path: str, prompt: str = "") -> str:
        self.emit("read_image", path)
        target = _resolve_safe_path(path, self.workspace_dir)
        if not target.exists():
            return f"Error: File not found: {path}"
        if not target.is_file():
            return f"Error: Not a file: {path}"
        header = target.read_bytes()[:32]
        media_type = self._image_media_type(target, header)
        if not media_type:
            return f"Error: Unsupported image type: {target.name}"
        image_bytes = target.read_bytes()
        if len(image_bytes) > 10 * 1024 * 1024:
            return f"Error: Image too large ({len(image_bytes)} bytes). Maximum size is 10MB."
        user_prompt = (prompt or "").strip() or "请描述这张图片的关键信息，并尽量识别其中可见文字。"
        reply = ""
        trace_id = f"img-{os.getpid()}-{target.name}"
        # #region debug-point V1:read-image-entry
        _debug_emit("V1", "tools.py:read_image", "read_image invoked", {
            "path": str(target),
            "media_type": media_type,
            "size": len(image_bytes),
            "model_id": self.model_id,
            "vision_model_id": self.vision_model_id or "",
            "has_vision_api_key": bool(self.vision_api_key),
            "has_anthropic_client": self.anthropic_client is not None,
        }, trace_id=trace_id)
        # #endregion
        if self._should_use_mmx_vision():
            reply = self._run_mmx_vision_describe(target, user_prompt, trace_id=trace_id)
            if reply:
                # #region debug-point V6:mmx-success
                _debug_emit("V6", "tools.py:read_image", "read_image succeeded via mmx", {
                    "reply_preview": reply[:300],
                }, trace_id=trace_id)
                # #endregion
                return reply
        if self.vision_api_key and self.vision_model_id:
            # #region debug-point V2:vision-branch
            _debug_emit("V2", "tools.py:read_image", "using dedicated vision backend", {
                "vision_model_id": self.vision_model_id,
                "vision_base_url": self.vision_base_url or "",
                "request_mode": "openai-chat-completions-image_url-data-uri",
                "prompt_len": len(user_prompt),
            }, trace_id=trace_id)
            # #endregion
            if self._vision_client is None:
                self._vision_client = OpenAI(
                    api_key=self.vision_api_key,
                    base_url=self.vision_base_url or None,
                )
            try:
                response = self._vision_client.chat.completions.create(
                    model=self.vision_model_id,
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {"type": "text", "text": user_prompt},
                                {
                                    "type": "image_url",
                                    "image_url": {
                                        "url": f"data:{media_type};base64,{base64.b64encode(image_bytes).decode('utf-8')}"
                                    },
                                },
                            ],
                        }
                    ],
                    max_tokens=1200,
                )
                # #region debug-point V7:vision-response
                _debug_emit("V7", "tools.py:read_image", "dedicated vision backend response summary", {
                    "model": getattr(response, "model", ""),
                    "finish_reason": getattr((response.choices or [None])[0], "finish_reason", ""),
                    "message_content_type": type(getattr(getattr((response.choices or [None])[0], "message", None), "content", "")).__name__,
                }, trace_id=trace_id)
                # #endregion
                reply = (response.choices[0].message.content or "").strip()
            except Exception as exc:
                # #region debug-point V2:vision-error
                _debug_emit("V2", "tools.py:read_image", "dedicated vision backend failed", {
                    "error": str(exc)[:500],
                }, trace_id=trace_id)
                # #endregion
                if _is_minimax_vision_model(self.vision_model_id or "") and "unknown model" in str(exc).lower():
                    return (
                        "Error: Current MiniMax key cannot access `MiniMax-VL-01` "
                        "(the API returned `unknown model`). "
                        "Please replace it with a MiniMax key that has VL-01 access."
                    )
                return f"Error: Vision tool failed: {exc}"
        elif _is_minimax_text_model(self.model_id):
            # #region debug-point V2:minimax-openai-compat
            _debug_emit("V2", "tools.py:read_image", "minimax text model cannot accept image input directly", {
                "model_id": self.model_id,
                "vision_model_id": self.vision_model_id or "",
                "has_api_key": bool(self.anthropic_api_key),
                "request_mode": "unsupported-minimax-text-image",
            }, trace_id=trace_id)
            # #endregion
            return (
                "Error: MiniMax text compatibility APIs do not support image input for `read_image`. "
                "Please configure a dedicated vision backend with "
                "`VISION_MODEL_ID=MiniMax-VL-01`. "
                "If `VISION_API_KEY` or `VISION_BASE_URL` is omitted, the tool will reuse the current MiniMax key "
                "and derive `https://api.minimaxi.com/v1` from `ANTHROPIC_BASE_URL`."
            )
        elif self.anthropic_client is not None:
            # #region debug-point V3:anthropic-branch
            _debug_emit("V3", "tools.py:read_image", "using anthropic image branch", {
                "model_id": self.model_id,
                "request_mode": "anthropic-messages-image-base64",
                "prompt_len": len(user_prompt),
            }, trace_id=trace_id)
            # #endregion
            try:
                response = self.anthropic_client.messages.create(
                    model=self.model_id,
                    max_tokens=1200,
                    messages=[
                        {
                            "role": "user",
                            "content": [
                                {
                                    "type": "image",
                                    "source": {
                                        "type": "base64",
                                        "media_type": media_type,
                                        "data": base64.b64encode(image_bytes).decode("utf-8"),
                                    },
                                },
                                {
                                    "type": "text",
                                    "text": user_prompt,
                                },
                            ],
                        }
                    ],
                )
                # #region debug-point V7:anthropic-response
                _debug_emit("V7", "tools.py:read_image", "anthropic vision response summary", {
                    "content_blocks": len(getattr(response, "content", []) or []),
                }, trace_id=trace_id)
                # #endregion
                texts = [getattr(block, "text", "") for block in getattr(response, "content", []) or [] if getattr(block, "text", "")]
                reply = "\n".join(texts).strip()
            except Exception as exc:
                # #region debug-point V3:anthropic-error
                _debug_emit("V3", "tools.py:read_image", "anthropic image branch failed", {
                    "error": str(exc)[:500],
                }, trace_id=trace_id)
                # #endregion
                return f"Error: Vision tool failed: {exc}"
        else:
            # #region debug-point V4:no-backend
            _debug_emit("V4", "tools.py:read_image", "no vision backend selected", {
                "model_id": self.model_id,
                "vision_model_id": self.vision_model_id or "",
                "has_vision_api_key": bool(self.vision_api_key),
                "anthropic_branch_blocked": self.anthropic_client is not None,
                "minimax_text_model": _is_minimax_text_model(self.model_id),
            }, trace_id=trace_id)
            # #endregion
            return (
                "Error: No vision backend configured. "
                "Set `VISION_API_KEY` and `VISION_MODEL_ID` "
                "(or `OPENAI_API_KEY` plus `VISION_MODEL_ID`) for the `read_image` tool. "
                "For MiniMax, use `VISION_MODEL_ID=MiniMax-VL-01`."
            )
        if not reply:
            # #region debug-point V5:empty-reply
            _debug_emit("V5", "tools.py:read_image", "vision model returned empty reply", {}, trace_id=trace_id)
            # #endregion
            return "Error: Vision model returned no text."
        # #region debug-point V6:success
        _debug_emit("V6", "tools.py:read_image", "read_image succeeded", {
            "reply_preview": reply[:300],
        }, trace_id=trace_id)
        # #endregion
        return reply

    def run_command(
        self,
        command: str,
        cwd: str = ".",
        timeout_seconds: int = _COMMAND_TIMEOUT_DEFAULT,
        max_output_chars: int = _COMMAND_OUTPUT_DEFAULT,
    ) -> str:
        normalized_command = str(command or "").strip()
        self.emit("run_command", normalized_command[:120])
        if not normalized_command:
            return "Error: command must not be empty."

        timeout = max(1, min(int(timeout_seconds), _COMMAND_TIMEOUT_MAX))
        output_limit = max(200, min(int(max_output_chars), _COMMAND_OUTPUT_MAX))
        supervision = self._supervise_command(
            normalized_command,
            cwd,
            timeout_seconds=timeout,
            max_output_chars=output_limit,
        )
        if supervision["status"] != "approved":
            return self._format_supervisor_rejection(supervision)

        resolved_cwd = supervision["cwd"]
        argv = supervision["argv"]
        return self._execute_command(
            command=normalized_command,
            argv=argv,
            resolved_cwd=resolved_cwd,
            timeout=timeout,
            output_limit=output_limit,
            risk=supervision["risk"],
            reason=supervision["reason"],
        )

    def get_pending_command(self, session_key: str = "") -> PendingCommandApproval | None:
        key = session_key or self._current_session_key
        if not key:
            return None
        return self._pending_command_approvals.get(key)

    def approve_pending_command(self, session_key: str = "", remember: bool = False) -> str:
        pending = self.get_pending_command(session_key)
        if pending is None:
            return "No pending command approval for the current session."
        if remember:
            prefix = self._approval_prefix(pending.argv)
            prefixes = self._approved_command_prefixes.setdefault(pending.session_key, [])
            if prefix not in prefixes:
                prefixes.append(prefix)
        self._pending_command_approvals.pop(pending.session_key, None)
        self._save_state()
        reason = pending.reason if not remember else f"{pending.reason}; remembered prefix {self._format_prefix(self._approval_prefix(pending.argv))}"
        return self._execute_command(
            command=pending.command,
            argv=pending.argv,
            resolved_cwd=Path(pending.cwd),
            timeout=pending.timeout_seconds,
            output_limit=pending.max_output_chars,
            risk=pending.risk,
            reason=self._append_persistence_warning(reason),
        )

    def reject_pending_command(self, session_key: str = "") -> str:
        pending = self.get_pending_command(session_key)
        if pending is None:
            return "No pending command approval for the current session."
        self._pending_command_approvals.pop(pending.session_key, None)
        self._save_state()
        self._record_supervisor_review(
            command=pending.command,
            argv=pending.argv,
            cwd=pending.cwd,
            risk=pending.risk,
            approved=False,
            reason="User rejected the pending command approval request.",
        )
        return "\n".join([
            "Supervisor: blocked (User rejected the pending command approval request.)",
            f"Risk: {pending.risk}",
            f"Command: {pending.command}",
            f"CWD: {pending.cwd}",
            self._state_warning_line(),
        ])

    def list_approved_command_prefixes(self, session_key: str = "") -> list[str]:
        key = session_key or self._current_session_key
        if not key:
            return []
        return [self._format_prefix(prefix) for prefix in self._approved_command_prefixes.get(key, [])]

    def format_pending_status_lines(self, session_key: str = "") -> list[str]:
        """Return pending command and remembered prefixes as display lines for CLI/Terminal."""
        key = session_key or self._current_session_key
        if not key:
            return []
        pending = self._pending_command_approvals.get(key)
        prefixes = self._approved_command_prefixes.get(key, [])
        lines = []
        if pending is None:
            lines.append("pending: (none)")
        else:
            lines.extend([
                f"pending: {pending.command}",
                f"risk: {pending.risk}",
                f"cwd: {pending.cwd}",
                f"requested_at: {pending.created_at}",
            ])
        if prefixes:
            lines.append("")
            lines.append("remembered prefixes:")
            lines.extend(f"- {self._format_prefix(p)}" for p in prefixes)
        else:
            lines.append("")
            lines.append("remembered prefixes: (none)")
        return lines

    def _execute_command(
        self,
        *,
        command: str,
        argv: list[str],
        resolved_cwd: Path,
        timeout: int,
        output_limit: int,
        risk: str,
        reason: str,
    ) -> str:
        env = os.environ.copy()
        env["PAGER"] = "cat"

        try:
            completed = subprocess.run(
                argv,
                cwd=str(resolved_cwd),
                capture_output=True,
                text=True,
                encoding="utf-8",
                errors="replace",
                timeout=timeout,
                env=env,
                shell=False,
            )
        except subprocess.TimeoutExpired as exc:
            partial = self._truncate_output((exc.stdout or "") + (exc.stderr or ""), output_limit)
            return "\n".join([
                f"Supervisor: approved ({reason})",
                f"Risk: {risk}",
                f"Command: {command}",
                f"CWD: {resolved_cwd}",
                f"Result: timed out after {timeout} second(s)",
                partial,
            ]).strip()
        except FileNotFoundError:
            return f"Error: Command not found: {argv[0]}"
        except Exception as exc:
            return f"Error: Failed to run command: {exc}"

        sections = [
            f"Supervisor: approved ({reason})",
            f"Risk: {risk}",
            f"Command: {command}",
            f"CWD: {resolved_cwd}",
            f"Exit code: {completed.returncode}",
        ]
        stdout_text = self._truncate_output(completed.stdout, output_limit)
        stderr_budget = max(0, output_limit - len(stdout_text))
        stderr_text = self._truncate_output(completed.stderr, stderr_budget) if stderr_budget else ""
        if stdout_text:
            sections.extend(["STDOUT:", stdout_text])
        if stderr_text:
            sections.extend(["STDERR:", stderr_text])
        if not stdout_text and not stderr_text:
            sections.append("(no output)")
        return "\n".join(sections)

    def _supervise_command(
        self,
        command: str,
        cwd: str,
        *,
        timeout_seconds: int,
        max_output_chars: int,
    ) -> dict[str, Any]:
        requested_cwd = cwd or "."
        try:
            resolved_cwd = _resolve_safe_path(requested_cwd, self.workspace_dir)
        except ValueError as exc:
            return self._reject_command(command, "high", str(exc), requested_cwd)
        if not resolved_cwd.exists():
            return self._reject_command(command, "high", f"Working directory not found: {requested_cwd}", requested_cwd)
        if not resolved_cwd.is_dir():
            return self._reject_command(command, "high", f"Working directory is not a directory: {requested_cwd}", requested_cwd)

        try:
            argv = shlex.split(command)
        except ValueError as exc:
            return self._reject_command(command, "high", f"Command parse failed: {exc}", str(resolved_cwd))
        if not argv:
            return self._reject_command(command, "high", "Command parse failed: empty argv", str(resolved_cwd))

        for marker in _SHELL_META_TOKENS:
            if marker in command:
                return self._reject_command(
                    command,
                    "high",
                    f"Shell operator `{marker}` is not allowed. Run a single command without pipes or redirects.",
                    str(resolved_cwd),
                    argv=argv,
                )

        executable = Path(argv[0]).name.lower()
        if executable in _ALWAYS_BLOCKED_COMMANDS:
            return self._reject_command(
                command,
                "high",
                f"Command `{executable}` is blocked by supervisor policy.",
                str(resolved_cwd),
                argv=argv,
            )

        risk = "medium"
        reason = "requires explicit user confirmation"
        if self._is_command_prefix_approved(argv):
            self._record_supervisor_review(
                command=command,
                argv=argv,
                cwd=str(resolved_cwd),
                risk=risk,
                approved=True,
                reason=f"matched remembered approval prefix {self._format_prefix(self._approval_prefix(argv))}",
            )
            return {
                "status": "approved",
                "risk": risk,
                "reason": f"matched remembered approval prefix {self._format_prefix(self._approval_prefix(argv))}",
                "argv": argv,
                "cwd": resolved_cwd,
            }

        if executable in {"pwd", "ls", "find", "grep", "rg", "cat", "head", "tail", "wc", "which", "whereis"}:
            risk = "low"
            reason = "read-only filesystem inspection"
        elif executable == "git":
            subcommand = argv[1].lower() if len(argv) > 1 else ""
            if subcommand in _BLOCKED_GIT_SUBCOMMANDS:
                return self._reject_command(
                    command,
                    "high",
                    f"`git {subcommand}` mutates repository state and is blocked.",
                    str(resolved_cwd),
                    argv=argv,
                )
            if subcommand in _MEDIUM_RISK_GIT_SUBCOMMANDS:
                risk = "low"
                reason = f"`git {subcommand}` is read-only"
        elif executable in {"python", "python3"}:
            if "-c" in argv:
                return self._reject_command(
                    command,
                    "high",
                    "Inline Python via `-c` is blocked. Put code in a file inside the workspace first.",
                    str(resolved_cwd),
                    argv=argv,
                )
            if any(part in {"-m", "pytest"} for part in argv[1:]) or len(argv) == 1 or argv[1] in {"-V", "--version"}:
                risk = "low"
                reason = "version check or test-style Python execution"
        elif executable == "pytest":
            risk = "low"
            reason = "test execution"
        elif executable in {"npm", "pnpm", "yarn"}:
            subcommand = argv[1].lower() if len(argv) > 1 else ""
            if subcommand in _BLOCKED_NODE_SUBCOMMANDS:
                return self._reject_command(
                    command,
                    "high",
                    f"`{executable} {subcommand}` may install or mutate dependencies and is blocked.",
                    str(resolved_cwd),
                    argv=argv,
                )
            if subcommand in _MEDIUM_RISK_NODE_SUBCOMMANDS:
                risk = "medium"
                reason = f"`{executable} {subcommand}` may execute project scripts"

        confirmed = risk == "low" or self._has_explicit_user_confirmation(command, argv)
        self._record_supervisor_review(
            command=command,
            argv=argv,
            cwd=str(resolved_cwd),
            risk=risk,
            approved=confirmed,
            reason=reason,
        )
        if not confirmed:
            return self._pend_command(
                command=command,
                argv=argv,
                cwd=str(resolved_cwd),
                risk=risk,
                reason=reason,
                timeout_seconds=timeout_seconds,
                max_output_chars=max_output_chars,
            )
        return {
            "status": "approved",
            "risk": risk,
            "reason": reason,
            "argv": argv,
            "cwd": resolved_cwd,
        }

    def _pend_command(
        self,
        *,
        command: str,
        argv: list[str],
        cwd: str,
        risk: str,
        reason: str,
        timeout_seconds: int,
        max_output_chars: int,
    ) -> dict[str, Any]:
        session_key = self._current_session_key or f"{self._current_channel}:{self._current_user_id}"
        self._pending_command_approvals[session_key] = PendingCommandApproval(
            session_key=session_key,
            command=command,
            argv=list(argv),
            cwd=cwd,
            risk=risk,
            reason=reason,
            timeout_seconds=timeout_seconds,
            max_output_chars=max_output_chars,
            created_at=datetime.now().isoformat(),
        )
        self._save_state()
        self._notify_pending_command(self._pending_command_approvals[session_key])
        return {
            "status": "pending",
            "risk": risk,
            "reason": self._append_persistence_warning(reason),
            "argv": argv,
            "cwd": cwd,
            "command": command,
            "session_key": session_key,
        }

    def _notify_pending_command(self, pending: PendingCommandApproval) -> None:
        if self._current_channel == "terminal" or self.message_sender is None:
            return
        try:
            self.message_sender(
                self._format_pending_command_message(pending),
                self._current_channel,
                self._current_user_id,
            )
        except Exception:
            pass

    def _format_pending_command_message(self, pending: PendingCommandApproval) -> str:
        lines = [
            "Supervisor requests approval",
            "",
            "Command:",
            pending.command,
            "",
            f"CWD: {pending.cwd}",
            f"Risk: {pending.risk}",
            f"Reason: {self._append_persistence_warning(pending.reason)}",
            "",
            "Reply with one of:",
            "/yes  allow once",
            "/always  allow and remember this command prefix",
            "/no  reject",
        ]
        warning_line = self._state_warning_line()
        if warning_line:
            lines.extend(["", warning_line])
        return "\n".join(lines)

    def _record_supervisor_review(
        self,
        *,
        command: str,
        argv: list[str],
        cwd: str,
        risk: str,
        approved: bool,
        reason: str,
    ) -> None:
        if self.supervisor_workflow is None:
            return
        self.supervisor_workflow.start_task(
            f"Run terminal command: {command}",
            requested_by=self._current_user_id,
            context={"channel": self._current_channel, "cwd": cwd},
        )
        self.supervisor_workflow.plan(
            objective="Evaluate whether a terminal command should be executed",
            steps=[
                "Parse the command without invoking a shell",
                "Check working directory access restrictions",
                "Classify risk and require confirmation when needed",
            ],
            affected_modules=["tools.py", "runtime"],
            involves_runtime=True,
            estimated_risk=risk,
        )
        report = self.supervisor_workflow.get_current_report()
        if report is not None and not approved:
            report.review_findings.append(
                ReviewFinding(
                    finding_id=f"cmd_{len(report.review_findings) + 1}",
                    phase=TaskPhase.REVIEW,
                    severity=Severity.BLOCKER,
                    rule_name="runtime_safety",
                    message=reason,
                    location=f"command:{Path(argv[0]).name if argv else command}",
                    recommendation="Request an explicit user confirmation or choose a safer read-only command.",
                )
            )
        summary = (
            f"runtime command precheck {'approved' if approved else 'blocked'}; "
            f"risk={risk}; reason={reason}; cwd={cwd}"
        )
        self.supervisor_workflow.execute(summary)
        self.supervisor_workflow.review()
        self.supervisor_workflow.finalize()

    def _reject_command(
        self,
        command: str,
        risk: str,
        reason: str,
        cwd: str,
        *,
        argv: list[str] | None = None,
        use_existing_report: bool = False,
    ) -> dict[str, Any]:
        if not use_existing_report:
            self._record_supervisor_review(
                command=command,
                argv=argv or [],
                cwd=cwd,
                risk=risk,
                approved=False,
                reason=reason,
            )
        return {
            "status": "blocked",
            "risk": risk,
            "reason": reason,
            "argv": argv or [],
            "cwd": cwd,
        }

    def _has_explicit_user_confirmation(self, command: str, argv: list[str]) -> bool:
        latest = self._latest_user_input.lower()
        if not latest:
            return False
        if not any(hint in latest for hint in _CONFIRMATION_HINTS):
            return False
        command_lower = command.lower()
        executable = Path(argv[0]).name.lower() if argv else ""
        return command_lower in latest or executable in latest

    def _approval_prefix(self, argv: list[str]) -> tuple[str, ...]:
        if not argv:
            return tuple()
        executable = Path(argv[0]).name.lower()
        if executable in {"git", "npm", "pnpm", "yarn", "python", "python3", "node"} and len(argv) >= 2:
            return tuple(part.lower() for part in argv[:2])
        return (executable,)

    def _is_command_prefix_approved(self, argv: list[str]) -> bool:
        if not self._current_session_key:
            return False
        prefix = self._approval_prefix(argv)
        return prefix in self._approved_command_prefixes.get(self._current_session_key, [])

    def _format_prefix(self, prefix: tuple[str, ...]) -> str:
        return " ".join(prefix) if prefix else "(empty)"

    def _format_supervisor_rejection(self, supervision: dict[str, Any]) -> str:
        status = supervision.get("status", "blocked")
        label = "pending approval" if status == "pending" else "blocked"
        lines = [f"Supervisor: {label} ({supervision['reason']})", f"Risk: {supervision['risk']}"]
        if supervision.get("command"):
            lines.append(f"Command: {supervision['command']}")
        lines.append(f"CWD: {supervision['cwd']}")
        warning_line = self._state_warning_line()
        if warning_line:
            lines.append(warning_line)
        if status == "pending" and self._current_channel == "terminal":
            lines.extend([
                "Next steps: use `/yes` to approve once, `/always` to approve and remember the command prefix, or `/no` to reject.",
            ])
        if self.supervisor_workflow is not None:
            report = self.supervisor_workflow.get_current_report()
            if report and report.gate_decision:
                lines.append(f"Decision: {report.gate_decision.overall.value}")
                lines.append(f"Reasoning: {report.gate_decision.reasoning}")
        return "\n".join(lines)

    def _truncate_output(self, text: str, max_chars: int) -> str:
        if max_chars <= 0:
            return ""
        normalized = text or ""
        if len(normalized) <= max_chars:
            return normalized
        return normalized[:max_chars] + f"\n\n[... truncated, total {len(normalized)} chars ...]"

    def _state_warning_line(self) -> str:
        if not self._state_error:
            return ""
        return f"Warning: {self._state_error}. Current approval still works in memory, but it may be lost after restart."

    def _append_persistence_warning(self, reason: str) -> str:
        warning_line = self._state_warning_line()
        if not warning_line:
            return reason
        return f"{reason}; {warning_line}"

    def format_prompt_block(self) -> str:
        lines = [
            "## Registered Tools",
            "",
            "Only call tools that appear in this section. Do not invent tool names.",
            "",
        ]
        for tool in self.tools:
            lines.append(f"- `{tool['name']}`: {tool['description']}")
        return "\n".join(lines)

    def process_tool_call(self, tool_name: str, tool_input: dict[str, Any]) -> str:
        handler = self.handlers.get(tool_name)
        if handler is None:
            return f"Error: Unknown tool '{tool_name}'"
        try:
            return handler(**tool_input)
        except TypeError as exc:
            return f"Error: Invalid arguments for {tool_name}: {exc}"
        except Exception as exc:
            return f"Error: {tool_name} failed: {exc}"
